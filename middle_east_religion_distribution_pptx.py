import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt

# 엑셀 파일 경로 수정
file_path = "C:/Users/sanghojeong9210/Desktop/코딩/middle_east_religion_distribution.xlsx"
df = pd.read_excel(file_path)

# 파워포인트 파일 생성
prs = Presentation()

# 슬라이드 레이아웃 설정 (제목 및 콘텐츠)
slide_layout = prs.slide_layouts[5]

for index, row in df.iterrows():
    country = row['국가']
    shia = row['시아파 (%)']
    sunni = row['수니파 (%)']
    christian = row['기독교 (%)']
    other = row['기타 (%)']
    
    # 슬라이드 추가
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = f"{country}의 종교 분포"
    
    # 차트 생성
    labels = ['시아파', '수니파', '기독교', '기타']
    sizes = [shia, sunni, christian, other]
    colors = ['#ff9999','#66b3ff','#99ff99','#ffcc99']
    fig1, ax1 = plt.subplots()
    ax1.pie(sizes, labels=labels, autopct='%1.1f%%', startangle=90, colors=colors)
    ax1.axis('equal')  # Equal aspect ratio ensures that pie is drawn as a circle.
    
    # 차트를 이미지로 저장
    chart_path = f"C:/Users/sanghojeong9210/Desktop/코딩/{country}_religion_distribution.png"
    plt.savefig(chart_path)
    plt.close(fig1)
    
    # 슬라이드에 이미지 삽입
    slide.shapes.add_picture(chart_path, Inches(1), Inches(1), width=Inches(8), height=Inches(5))

# 파워포인트 파일 저장 경로 수정
pptx_path = "C:/Users/sanghojeong9210/Desktop/코딩/middle_east_religion_distribution.pptx"
prs.save(pptx_path)
